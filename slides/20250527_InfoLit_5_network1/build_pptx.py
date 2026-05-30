"""Build a fully editable native PPTX deck from the lecture data.

All shapes (header band, page-title trapezoid, logo, body text, takeaway,
section boxes, etc.) are real PowerPoint shapes — every element is editable.

Run:
    python3 build_pptx.py
Output:
    out/20250527_InfoLit_5_network1_editable.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

DECK_DIR = Path(__file__).parent
SRC = DECK_DIR / "src"
LOGO = DECK_DIR.parent / "assets" / "chiba-logo.png"
OUT = DECK_DIR / "out" / "20250527_InfoLit_5_network1_editable.pptx"

# --- design tokens -----------------------------------------------------------
ACCENT = RGBColor(0xA6, 0x19, 0x2E)       # 千葉大 red
ACCENT_DARK = RGBColor(0x7E, 0x10, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xFA, 0xF6, 0xF7)

DECK_TITLE = "情報リテラシー(5) ネットワーク(1)"
FONT = "Hiragino Kaku Gothic ProN"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HDR_H = Inches(0.55)

# ---------------------------------------------------------------------------
def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()

def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=GRAY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, bg=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    if bg is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_bullets(slide, left, top, width, height, items, *, size=16, color=GRAY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "・" + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb

def header_band(slide, page_title):
    # Left red band (deck title)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.6), HDR_H)
    set_solid(band, ACCENT)
    tf = band.text_frame
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = DECK_TITLE
    r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

    # Center darker red page-title block
    pt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), 0, Inches(5.4), HDR_H)
    set_solid(pt, ACCENT_DARK)
    tf = pt.text_frame
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = page_title
    r.font.name = FONT; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE

    # Thin red bottom line on right
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.1), HDR_H - Inches(0.05), Inches(3.233), Inches(0.05))
    set_solid(line, ACCENT)

    # Logo top-right
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.6), Inches(0.07), height=Inches(0.42))

def add_subtitle(slide, text, *, top=Inches(0.75)):
    add_text(slide, Inches(0.5), top, Inches(12.3), Inches(0.55),
             text, size=24, bold=True, color=ACCENT_DARK)

def add_takeaway(slide, text):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.55))
    set_solid(box, ACCENT)
    tf = box.text_frame
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE

def add_page_no(slide, n):
    add_text(slide, Inches(12.7), Inches(7.15), Inches(0.6), Inches(0.3),
             str(n), size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def fit_picture(slide, img_path, box_l, box_t, box_w, box_h):
    """Add picture, centered within (box_l, box_t, box_w, box_h)."""
    from PIL import Image as PILImage
    im = PILImage.open(img_path)
    iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    l = int(box_l + (box_w - w) / 2)
    t = int(box_t + (box_h - h) / 2)
    slide.shapes.add_picture(str(img_path), l, t, w, h)

# --- slide builders ---------------------------------------------------------
def slide_cover(prs, *, title, subtitle, author, affil, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Top accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.4))
    set_solid(bar, ACCENT)
    # Title
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.4),
             title, size=44, bold=True, color=ACCENT_DARK)
    add_text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.8),
             subtitle, size=24, color=GRAY)
    # Author block
    add_text(s, Inches(0.8), Inches(5.2), Inches(8), Inches(0.5),
             author, size=20, bold=True, color=GRAY)
    add_text(s, Inches(0.8), Inches(5.7), Inches(8), Inches(0.5),
             affil, size=16, color=GRAY)
    add_text(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
             meta, size=14, color=GRAY)
    # Logo bottom-right
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(11.0), Inches(6.4), height=Inches(0.6))
    return s

def slide_fig(prs, n, *, page_title, subtitle, img, takeaway):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_band(s, page_title)
    add_subtitle(s, subtitle)
    fit_picture(s, SRC / img, Inches(0.7), Inches(1.4),
                Inches(11.9), Inches(5.2))
    add_takeaway(s, takeaway)
    add_page_no(s, n)
    return s

def slide_split(prs, n, *, page_title, subtitle, img, heading, bullets, takeaway):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_band(s, page_title)
    add_subtitle(s, subtitle)
    # Left: figure
    fit_picture(s, SRC / img, Inches(0.5), Inches(1.4),
                Inches(6.5), Inches(5.2))
    # Right: heading + bullets
    add_text(s, Inches(7.3), Inches(1.45), Inches(5.5), Inches(0.5),
             heading, size=20, bold=True, color=ACCENT_DARK)
    add_bullets(s, Inches(7.3), Inches(2.0), Inches(5.7), Inches(4.6),
                bullets, size=15)
    add_takeaway(s, takeaway)
    add_page_no(s, n)
    return s

def slide_summary(prs, n, *, page_title, subtitle, sections):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_band(s, page_title)
    add_subtitle(s, subtitle)
    # 2 section boxes side-by-side
    cols = len(sections)
    gap = Inches(0.3)
    total_w = Inches(12.3)
    box_w = (total_w - gap * (cols - 1)) / cols
    top = Inches(1.5)
    box_h = Inches(5.0)
    for i, (head, items) in enumerate(sections):
        left = Inches(0.5) + (box_w + gap) * i
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = ACCENT; box.line.width = Pt(1.5)
        # Title strip
        add_text(s, left, top + Inches(0.15), box_w, Inches(0.6),
                 head, size=18, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
        add_bullets(s, left + Inches(0.15), top + Inches(0.9),
                    box_w - Inches(0.3), box_h - Inches(1.0),
                    items, size=15)
    add_page_no(s, n)
    return s

def slide_divider(prs, n, *, chapter, title, sub):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_solid(bg, ACCENT)
    add_text(s, Inches(1.0), Inches(2.6), Inches(10), Inches(0.6),
             f"CHAPTER {chapter}", size=22, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(3.2), Inches(11), Inches(1.4),
             title, size=64, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(4.8), Inches(11), Inches(0.7),
             sub, size=22, color=WHITE)
    add_page_no(s, n)
    return s

def slide_refs(prs, n, *, refs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header_band(s, "参考文献")
    add_subtitle(s, "References")
    add_bullets(s, Inches(0.7), Inches(1.6), Inches(12), Inches(5.5),
                refs, size=15)
    add_page_no(s, n)
    return s

def slide_qa(prs, n, *, contact):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_solid(bg, ACCENT)
    add_text(s, Inches(0), Inches(2.6), SLIDE_W, Inches(1.8),
             "Q&A", size=120, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.7), SLIDE_W, Inches(0.6),
             contact, size=22, color=WHITE, align=PP_ALIGN.CENTER)
    add_page_no(s, n)
    return s

# --- deck content -----------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    slide_cover(prs,
        title="情報リテラシー 講義(5)",
        subtitle="ネットワークの仕組みとインターネット (1)",
        author="檜垣 泰彦",
        affil="千葉大学 アカデミック・リンク・センター 特任教授",
        meta="2025年5月27日 / 千葉大学Moodle 情報リテラシー(23)(24) T1-2・火4")

    # 2. LAN
    slide_fig(prs, 2,
        page_title="LAN (Local Area Network)",
        subtitle="1つの建物・キャンパスでのネットワーク",
        img="slide02.png",
        takeaway="スイッチ・ハブで端末を束ね、限られた地域内で相互接続する")

    # 3. WAN
    slide_fig(prs, 3,
        page_title="WAN (Wide Area Network)",
        subtitle="離れた地域のLANを接続する",
        img="slide03.png",
        takeaway="ルーターを介してLAN同士を結び、地理的に分散した通信を実現")

    # 4. ネットワーク構成要素
    slide_fig(prs, 4,
        page_title="ネットワークの構成要素",
        subtitle="実環境を構成する装置群",
        img="slide04.png",
        takeaway="ルーター/L2・L3スイッチ/FW/無線APなど用途別の機器が連携")

    # 5. インターネットの構造
    slide_split(prs, 5,
        page_title="インターネットの構造",
        subtitle="ISP・IX・NOC",
        img="slide05.png",
        heading="主要な略語",
        bullets=[
            "ISP：Internet Service Provider",
            "IX：Internet Exchange",
            "NOC：Network Operation Center",
            "イントラネット：組織内ネットワーク",
            "輻輳（ふくそう）：通信が許容量を超え繋がりにくくなる現象",
        ],
        takeaway="TCP/IPで結ばれたネットワークの集合体 ＝ インターネット")

    # 6. 言語層と装置層
    slide_split(prs, 6,
        page_title="言語層と装置層の2階層モデル",
        subtitle="プロトコル：通信規約",
        img="slide06.png",
        heading="プロトコルとは",
        bullets=[
            "ネットワーク通信のための約束事の集合",
            "メーカや OS が異なる機器でも互いに通信可能",
            "言語層と装置層は独立：片方だけ差し替え可能",
        ],
        takeaway="層に分けて考えることで「片方だけ変更」が可能になる")

    # 7. OSI 7階層
    slide_fig(prs, 7,
        page_title="OSI参照モデル（7階層）",
        subtitle="層に分けて考える（OSI: Open System Interconnection）",
        img="slide07.png",
        takeaway="アプリケーション層〜物理層の7層構造。ルーターはL3まで処理")

    # 8. ブリッジ・ハブ
    slide_fig(prs, 8,
        page_title="データリンク層の装置【ブリッジ・ハブ】",
        subtitle="L2機器の役割",
        img="slide08.png",
        takeaway="スイッチはポートごとのブリッジ。PoEで給電も可能")

    # 9. ルーター
    slide_fig(prs, 9,
        page_title="ネットワーク層の装置【ルーター】",
        subtitle="経路を決めてパケットを配送する配達屋",
        img="slide09.png",
        takeaway="宛先へのルートを決定し、任意のデータリンク同士を接続できる")

    # 10. 各機器と階層
    slide_fig(prs, 10,
        page_title="各機器と該当する階層",
        subtitle="リピーター／ブリッジ／ルーター",
        img="slide10.png",
        takeaway="物理層=リピーター／L2=ブリッジ・L2SW／L3=ルーター・L3SW")

    # 11. TCP/IP群
    slide_split(prs, 11,
        page_title="TCP/IPプロトコル群",
        subtitle="もっとも有名でよく使われているプロトコル",
        img="slide11.png",
        heading="4種類のプロトコル",
        bullets=[
            "アプリケーション：HTTP, SMTP, FTP, TELNET, SNMP",
            "トランスポート：TCP, UDP",
            "インターネット：IP, ICMP, ARP",
            "経路制御：RIP, OSPF, BGP",
        ],
        takeaway="TCP/IP ⇔ インターネット。事実上の標準プロトコル群")

    # 12. OSIとTCP/IP対応
    slide_fig(prs, 12,
        page_title="OSI参照モデルとTCP/IPの関係",
        subtitle="7階層モデルと4階層モデルの対応",
        img="slide12.png",
        takeaway="OSIの5-7層 ≒ TCP/IPアプリケーション層。物理層はハードウェアが担う")

    # 13. Divider CH1
    slide_divider(prs, 13,
        chapter=1, title="データリンク層",
        sub="MACアドレスとイーサネット／無線LAN")

    # 14. MAC
    slide_split(prs, 14,
        page_title="データリンクの技術",
        subtitle="MACアドレスで宛先を判定",
        img="slide14.png",
        heading="MACアドレス",
        bullets=[
            "Media Access Control Address",
            "機器に一意に割り当てられた物理アドレス",
            "例：E0:69:95:62:D5:F1",
            "確認(Win11)：コントロールパネル → ネットワークと共有 → 詳細",
        ],
        takeaway="L2では IP ではなく MAC アドレスでフレームを届ける")

    # 15. イーサネット種類
    slide_split(prs, 15,
        page_title="イーサネットの種類",
        subtitle="速度とケーブル規格",
        img="slide15.png",
        heading="主な規格（最大長は概ね 100m）",
        bullets=[
            "10BASE-T：10 Mbps / カテゴリ 3〜5",
            "100BASE-TX：100 Mbps / カテゴリ 5",
            "1000BASE-T：1 Gbps / カテゴリ 5/5e",
            "10GBASE-T：10 Gbps / カテゴリ 6a",
        ],
        takeaway="対応カテゴリのケーブルを使わないと速度が出ない")

    # 16. 802.11
    slide_fig(prs, 16,
        page_title="IEEE802.11（無線通信）の比較",
        subtitle="物理層・MAC層の規格",
        img="slide16.png",
        takeaway="11n=150Mbps、11ac=6.93Gbps、11ax=9.6Gbps と高速化が進む")

    # 17. WiFi世代
    slide_fig(prs, 17,
        page_title="最近のWiFi規格",
        subtitle="Wi-Fi 4 〜 Wi-Fi 7",
        img="slide17.png",
        takeaway="Wi-Fi 7（2024年, 11be）は36Gbps・6GHz帯・320MHz帯域幅に対応")

    # 18. チャンネルレイアウト
    slide_fig(prs, 18,
        page_title="Wi-Fiにおけるチャンネルレイアウト",
        subtitle="2.4 / 5 / 6 GHz帯と日本の利用可能帯域",
        img="slide18.png",
        takeaway="6GHz帯（480MHz幅）まで拡張、Wi-Fi 6E/7で活用")

    # 19. 802.11aチャンネル
    slide_fig(prs, 19,
        page_title="802.11a のチャンネル",
        subtitle="5GHz帯と気象レーダー干渉",
        img="slide19.png",
        takeaway="W56（2007年〜）から屋外利用が解禁、ただし気象レーダー干渉あり")

    # 20. inSSIDer
    slide_fig(prs, 20,
        page_title="学内における無線LANの電波状況",
        subtitle="inSSIDer による観測",
        img="slide20.png",
        takeaway="2.4GHz帯はチャネルが密集し干渉、5GHzは比較的空いている")

    # 21. 無線LAN用語
    slide_fig(prs, 21,
        page_title="無線LANの用語",
        subtitle="識別子・暗号・認証",
        img="slide21.png",
        takeaway="SSID / WEP・WPA・WPA2 / PSK / IEEE802.1X / EAP / WPS")

    # 22. セキュリティ
    slide_summary(prs, 22,
        page_title="無線LANのセキュリティ",
        subtitle="暗号方式と接続制御",
        sections=[
            ("暗号方式の優先順", [
                "WPA2-PSK (AES) を使うのが一番安全",
                "非対応なら WPA-PSK (AES)",
                "WEPの解読は非常に簡単：使ってはいけない",
            ]),
            ("補助的な対策", [
                "MACアドレスフィルタリング",
                "ESSIDステルス機能で発信停止",
                "any接続拒否でPC側のESSID一覧に出さない",
            ]),
        ])

    # 23. ADSL
    slide_fig(prs, 23,
        page_title="ADSL接続",
        subtitle="アナログ電話回線でのデジタル通信",
        img="slide23.png",
        takeaway="非対称デジタル加入者線、1.5〜50Mbps。2023〜2024年にサービス終了")

    # 24. FTTH
    slide_fig(prs, 24,
        page_title="FTTH接続",
        subtitle="光ファイバを家庭まで",
        img="slide24.png",
        takeaway="ONU（光回線終端装置）と光スプリッタで各家庭まで光配信")

    # 25. CATV
    slide_fig(prs, 25,
        page_title="ケーブルテレビによるインターネット接続",
        subtitle="ヘッドエンド経由でISPへ",
        img="slide25.png",
        takeaway="CATVの同軸/光網を流用、ケーブルモデムでISPへ接続")

    # 26. VPN
    slide_split(prs, 26,
        page_title="VPN (Virtual Private Network)",
        subtitle="公衆網上の仮想専用線",
        img="slide26.png",
        heading="VPNとは",
        bullets=[
            "公衆回線を暗号化・認証で安全化",
            "あたかも専用線で結んだように利用",
            "モバイル環境から組織内へ安全接続にも",
            "ラベル付け（MPLS）で仮想的な独自網を構築",
        ],
        takeaway="公衆網を「専用線化」して安全に組織ネットワークへ接続")

    # 27. Divider CH2
    slide_divider(prs, 27,
        chapter=2, title="IPプロトコル",
        sub="ネットワーク層のプロトコル")

    # 28. IPの役割
    slide_fig(prs, 28,
        page_title="IPの役割",
        subtitle="複雑なネットワークの中の宛先まで届ける",
        img="slide28.png",
        takeaway="複数ルーターを経由しても最終ホストまでパケットを届けるのがIP")

    # 29. IPアドレス
    slide_fig(prs, 29,
        page_title="IPアドレス",
        subtitle="インターネット上のホスト識別子",
        img="slide29.png",
        takeaway="グローバルIPはICANNにより世界で唯一の値を割当")

    # 30. NIC毎にIP
    slide_split(prs, 30,
        page_title="NIC毎に1つ以上のIPアドレス",
        subtitle="サブネットマスクで分ける",
        img="slide30.png",
        heading="ポイント",
        bullets=[
            "NIC：Network Interface Card",
            "ホストには最低1つのIPアドレス",
            "ルーターは2つ以上（境界となるため）",
            "1つのNICに複数IPも可",
            "サブネットマスクで NW部・ホスト部 を区切る",
        ],
        takeaway="192.168.200.107/24 の /24 はネットワーク部のビット数")

    # 31. ホスト部
    slide_fig(prs, 31,
        page_title="IPアドレスのホスト部",
        subtitle="同一セグメント内での約束事",
        img="slide31.png",
        takeaway="同一セグメント＝ネットワーク部が同じ、ホスト部は重複不可")

    # 32. ネットワーク部
    slide_fig(prs, 32,
        page_title="IPアドレスのネットワーク部",
        subtitle="ルーターによる経路判断",
        img="slide32.png",
        takeaway="ルーターは宛先IPのネットワーク部を見て転送先を決定する")

    # 33. ホスト名
    slide_fig(prs, 33,
        page_title="ホスト名とIPアドレス",
        subtitle="わかりやすい名前と覚えにくいアドレスの対応",
        img="slide33.png",
        takeaway="hostsファイル：ホスト名⇔IPアドレスの対応表（DNS以前の仕組み）")

    # 34. ドメイン階層
    slide_fig(prs, 34,
        page_title="ドメインの階層構造",
        subtitle="例：tu.chiba-u.ac.jp",
        img="slide34.png",
        takeaway="root → jp → ac → chiba-u → tu と右から左へ階層を辿る")

    # 35. DNS
    slide_fig(prs, 35,
        page_title="ネームサーバ（DNS）",
        subtitle="名前情報の分散型データベース",
        img="slide35.png",
        takeaway="ルート → トップレベル → 各組織DNS へと再帰的に問い合わせ")

    # 36. DHCP
    slide_fig(prs, 36,
        page_title="DHCP (Dynamic Host Configuration Protocol)",
        subtitle="IPアドレスを自動配布する仕組み",
        img="slide36.png",
        takeaway="接続するだけでIP・サブネット・GW・DNSが設定される")

    # 37. NAT
    slide_fig(prs, 37,
        page_title="NAT (Network Address Translation)",
        subtitle="プライベートIP ⇔ グローバルIP の変換",
        img="slide37.png",
        takeaway="家庭ルータが備える機能。10/8・172.16/12・192.168/16 を変換")

    # 38. FW
    slide_fig(prs, 38,
        page_title="ファイアウォール（防火壁）",
        subtitle="内部と外部の境界に置く防御装置",
        img="slide38.png",
        takeaway="DMZ（非武装地帯）に公開サーバ、内部への侵入は禁止する")

    # 39. IPv6
    slide_summary(prs, 39,
        page_title="IPv6 (IP version 6)",
        subtitle="IPv4アドレス枯渇への対策",
        sections=[
            ("基本", [
                "IPv4は 32ビット、IPv6は 128ビット",
                "直接アクセス可能なアドレス不足が解消",
                "IPv4における問題の解決を目指す",
            ]),
            ("IPv6の特徴", [
                "IPアドレス拡大と経路制御表の集約",
                "パフォーマンスの向上",
                "プラグ＆プレイ機能を必須化",
                "認証機能・暗号化機能の採用",
            ]),
        ])

    # 40. Divider CH3
    slide_divider(prs, 40,
        chapter=3, title="TCPとUDP",
        sub="トランスポート層のプロトコル")

    # 41. ポート
    slide_fig(prs, 41,
        page_title="ポート番号によるアプリケーションの識別",
        subtitle="同一ホスト内で複数サービスを識別",
        img="slide41.png",
        takeaway="IPで「ホスト」、ポート番号で「アプリ」を識別する")

    # 42. TCP wellknown
    slide_fig(prs, 42,
        page_title="TCPのウェルノウンポート（一部）",
        subtitle="主要サービスの予約番号",
        img="slide42.png",
        takeaway="20/21=FTP、22=SSH、25=SMTP、53=DNS、80=HTTP を押さえる")

    # 43. UDP wellknown
    slide_fig(prs, 43,
        page_title="UDPのウェルノウンポート（一部）",
        subtitle="DNS・DHCP・NTPなどはUDP",
        img="slide43.png",
        takeaway="53=DNS、67/68=DHCP、69=TFTP、123=NTP、161=SNMP")

    # 44. Refs
    slide_refs(prs, 44, refs=[
        "竹下 隆史, 村山 公保, 荒井 透, 苅田 幸雄 (2012). マスタリングTCP/IP 入門編 第5版. オーム社.",
        "竹下 隆史, 村山 公保, 荒井 透, 苅田 幸雄 (2019). マスタリングTCP/IP 入門編 第6版. オーム社.",
        "KEYENCE (2024). Wi-Fi 規格と最新動向. https://www.keyence.co.jp/ss/products/controls/wifi/basis/latest_standards.jsp",
        "平成25年春期 ITパスポート試験公開問題 問68.",
    ])

    # 45. QA
    slide_qa(prs, 45, contact="higaki.yasuhiko@faculty.chiba-u.jp")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"saved: {OUT}  ({OUT.stat().st_size/1024/1024:.1f} MB)")

if __name__ == "__main__":
    build()
