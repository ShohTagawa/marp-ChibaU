#!/usr/bin/env python3
"""
pptx2marp — PowerPoint(.pptx) を「座標再現＋academic CSS」方式で Marp デックに変換する。

設計方針（ユーザー要件）:
  - ヘッダー以外は **元スライドと完全一致**：各図形/画像/テキストを元の EMU 座標・サイズの
    まま % に変換して HTML で絶対配置する（位置・サイズ・順序・内容を保持）。
  - **デザインだけ academic を優先**：配色・字体（font-family）・ヘッダー帯・ページ番号は
    academic テーマ（theme/academic.css）に任せる。色や書体はインライン指定せず CSS に勝たせる。
  - フォント「サイズ」だけは元 pt を維持（箱からの溢れ防止＝レイアウト一致のため）。
  - 各スライドのタイトル枠は academic の中央帯 `page-title` に載せる（＝ヘッダー扱い）。

出力（CLAUDE.md 規約）:
  slides/<deck-name>/
    <deck-name>.md
    src/      (抽出画像 figNN-...)
    out/

依存: python-pptx  (pip3 install python-pptx)
使い方:
  python3 tools/pptx2marp/pptx2marp.py <input.pptx> [--name ...] [--title ...] [--event ...] [--date ...]
"""
from __future__ import annotations

import argparse
import datetime as _dt
import html
import re
import shutil
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.oxml.ns import qn
except ImportError:
    sys.exit("python-pptx が必要です:  pip3 install python-pptx")


REPO = Path(__file__).resolve().parents[2]
THEME_CSS = REPO / "theme" / "academic.css"

UNSUPPORTED_IMG = {"wmf", "emf"}  # Chrome/Marp が PDF 化で空白にするベクター形式
BODY_PH = {
    PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT,
    getattr(PP_PLACEHOLDER, "CONTENT", PP_PLACEHOLDER.OBJECT),
}
SUBTITLE_PH = {PP_PLACEHOLDER.SUBTITLE}


def is_title_slide(s) -> bool:
    name = (s.slide_layout.name or "")
    if "タイトル スライド" in name or "Title Slide" in name:
        return True
    for sh in s.shapes:
        if sh.is_placeholder and sh.placeholder_format.type in SUBTITLE_PH:
            return True
    return False
EMU_PER_PT = 12700
MARP_H_PX = 720.0  # Marp の section 高さは 16:9 / 4:3 とも 720px


# --------------------------------------------------------------------------- #
# helpers
# --------------------------------------------------------------------------- #
def slugify(text: str) -> str:
    text = re.sub(r"[^A-Za-z0-9\- ]+", "", text)
    text = re.sub(r"\s+", "-", text.strip())
    return re.sub(r"-{2,}", "-", text).strip("-").lower()


def esc(s: str) -> str:
    return html.escape(s, quote=True)


def text_to_html(s: str) -> str:
    """段落内テキスト。改行(ソフト/ハード)は <br> として保持。"""
    s = s.replace("\r", "").replace("\x0b", "\n").replace("\v", "\n").replace("\t", "　")
    return "<br>".join(esc(line) for line in s.split("\n"))


def get_image(shape):
    try:
        return shape.image
    except Exception:
        return None


ALIGN_MAP = {
    "CENTER (2)": "center", "RIGHT (3)": "right",
    "JUSTIFY (4)": "justify", "JUSTIFY_LOW (7)": "justify",
}


def para_align(para):
    return ALIGN_MAP.get(str(para.alignment), None)


def para_font_pt(para):
    if para.font.size is not None:
        return para.font.size.pt
    for r in para.runs:
        if r.font.size is not None:
            return r.font.size.pt
    return None


def para_bullet(para, is_body: bool):
    """段落の箇条書き記号を返す。なければ None。"""
    pPr = para._p.find(qn("a:pPr"))
    if pPr is not None:
        if pPr.find(qn("a:buNone")) is not None:
            return None
        bu = pPr.find(qn("a:buChar"))
        if bu is not None:
            return bu.get("char") or "•"
        if pPr.find(qn("a:buAutoNum")) is not None:
            return "num"
    return "•" if is_body else None


# --------------------------------------------------------------------------- #
# geometry / group transform
# --------------------------------------------------------------------------- #
def make_transform(parent_tf, group_shape):
    """グループの子座標系 → 親座標系 への変換関数を合成して返す。"""
    el = group_shape._element
    xfrm = None
    grpSpPr = el.find(qn("p:grpSpPr"))
    if grpSpPr is not None:
        xfrm = grpSpPr.find(qn("a:xfrm"))
    if xfrm is None:
        return parent_tf
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    chOff, chExt = xfrm.find(qn("a:chOff")), xfrm.find(qn("a:chExt"))
    if None in (off, ext, chOff, chExt):
        return parent_tf
    gx, gy = int(off.get("x")), int(off.get("y"))
    gw, gh = int(ext.get("cx")), int(ext.get("cy"))
    cx, cy = int(chOff.get("x")), int(chOff.get("y"))
    cw, ch = int(chExt.get("cx")) or 1, int(chExt.get("cy")) or 1
    sx, sy = gw / cw, gh / ch

    def tf(x, y, w, h):
        X = gx + (x - cx) * sx
        Y = gy + (y - cy) * sy
        return parent_tf(X, Y, w * sx, h * sy)

    return tf


# --------------------------------------------------------------------------- #
# slide walk
# --------------------------------------------------------------------------- #
class Slide:
    def __init__(self):
        self.title = ""
        self.elements = []   # dict: kind, geo(l,t,w,h %), payload...
        self.notes = ""


def geo_pct(transform, shape, W, H):
    l = shape.left if shape.left is not None else int(0.05 * W)
    t = shape.top if shape.top is not None else int(0.16 * H)
    w = shape.width if shape.width is not None else int(0.90 * W)
    h = shape.height if shape.height is not None else int(0.10 * H)
    X, Y, Wd, Hd = transform(l, t, w, h)
    return {
        "l": round(X / W * 100, 3), "t": round(Y / H * 100, 3),
        "w": round(Wd / W * 100, 3), "h": round(Hd / H * 100, 3),
    }


def walk(shapes, transform, slide, ctx, title_shape, lift_title):
    W, H, src_dir, fig, scale = ctx["W"], ctx["H"], ctx["src"], ctx["fig"], ctx["scale"]
    for shape in shapes:  # 文書順＝z順を維持（重なりを正しく）
        # グループ → 再帰
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(shape.shapes, make_transform(transform, shape), slide, ctx, title_shape, lift_title)
            continue

        # タイトル枠：通常スライドは academic の page-title に載せる（ヘッダー扱い）。
        # 表紙レイアウトは元座標のまま大見出しとして置く（lift_title=False）。
        if lift_title and title_shape is not None and shape == title_shape:
            if shape.has_text_frame:
                slide.title = " ".join(shape.text_frame.text.split())
            continue

        # 画像
        img = get_image(shape)
        if img is not None:
            fig[0] += 1
            ext = (img.ext or "png").lower()
            ext = "jpg" if ext == "jpeg" else ext
            fname = f"fig{fig[0]:02d}-img.{ext}"
            (src_dir / fname).write_bytes(img.blob)
            g = geo_pct(transform, shape, W, H)
            slide.elements.append({
                "kind": "image", "geo": g,
                "src": f"./src/{fname}", "ok": ext not in UNSUPPORTED_IMG,
            })
            continue

        # 表 → HTML テーブルで座標再現
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for r in shape.table.rows:
                rows.append([text_to_html(c.text) for c in r.cells])
            slide.elements.append({
                "kind": "table", "geo": geo_pct(transform, shape, W, H), "rows": rows,
            })
            continue

        # グラフ等の GraphicFrame（埋め込み再現不可）→ 位置に注記
        if getattr(shape, "has_chart", False) and shape.has_chart:
            slide.elements.append({
                "kind": "ph", "geo": geo_pct(transform, shape, W, H),
                "label": "📊 グラフ（元ファイル参照 / marp-echarts で要再作成）",
            })
            continue

        # 動画（埋め込み不可 → 位置にプレースホルダ）
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            slide.elements.append({
                "kind": "ph", "geo": geo_pct(transform, shape, W, H),
                "label": f"🎬 動画: {esc(shape.name or 'media')}（元ファイル参照）",
            })
            continue

        # テキスト
        if shape.has_text_frame and shape.text_frame.text.strip():
            is_body = shape.is_placeholder and shape.placeholder_format.type in BODY_PH
            paras = []
            autonum = 0
            for para in shape.text_frame.paragraphs:
                raw = para.text
                if not raw.strip():
                    paras.append({"html": "&#8203;", "size": None, "align": None,
                                  "level": para.level or 0, "bullet": None})
                    continue
                bullet = para_bullet(para, is_body)
                if bullet == "num":
                    autonum += 1
                    bullet = f"{autonum}."
                pt = para_font_pt(para)
                paras.append({
                    "html": text_to_html(raw),
                    "size": round(pt * scale, 1) if pt else None,
                    "align": para_align(para),
                    "level": para.level or 0,
                    "bullet": bullet,
                })
            slide.elements.append({
                "kind": "text", "geo": geo_pct(transform, shape, W, H), "paras": paras,
            })
    return slide


# --------------------------------------------------------------------------- #
# markdown emission
# --------------------------------------------------------------------------- #
def render_text_box(el) -> str:
    g = el["geo"]
    style = (f"left:{g['l']}%;top:{g['t']}%;width:{g['w']}%;height:{g['h']}%")
    lines = [f'<div class="ppt-box" style="{style}">']
    for p in el["paras"]:
        st = []
        if p["size"]:
            st.append(f"font-size:{p['size']}px")
        if p["align"]:
            st.append(f"text-align:{p['align']}")
        if p["level"]:
            st.append(f"margin-left:{p['level'] * 1.3:.1f}em")
        sattr = f' style="{";".join(st)}"' if st else ""
        bullet = f'<span class="ppt-bu">{esc(p["bullet"])}</span>' if p["bullet"] else ""
        lines.append(f'<div class="ppt-p"{sattr}>{bullet}{p["html"]}</div>')
    lines.append("</div>")
    return "\n".join(lines)


def render_image(el) -> str:
    g = el["geo"]
    style = f"left:{g['l']}%;top:{g['t']}%;width:{g['w']}%;height:{g['h']}%"
    if el["ok"]:
        return f'<img class="ppt-img" style="{style}" src="{el["src"]}" />'
    name = Path(el["src"]).name
    return (f'<div class="ppt-ph" style="{style}">[ベクター画像 {esc(name)}<br>'
            f'PowerPointでPNG等に再保存してください]</div>')


def render_table(el) -> str:
    g = el["geo"]
    style = f"left:{g['l']}%;top:{g['t']}%;width:{g['w']}%;height:{g['h']}%"
    lines = [f'<table class="ppt-tbl" style="{style}">']
    for i, row in enumerate(el["rows"]):
        lines.append("<tr>")
        tag = "th" if i == 0 else "td"
        for cell in row:
            lines.append(f"<{tag}>{cell}</{tag}>")
        lines.append("</tr>")
    lines.append("</table>")
    return "\n".join(lines)


def render_ph(el) -> str:
    g = el["geo"]
    style = f"left:{g['l']}%;top:{g['t']}%;width:{g['w']}%;height:{g['h']}%"
    return f'<div class="ppt-ph" style="{style}">{el["label"]}</div>'


def render_slide(slide: Slide) -> str:
    out = ["<!-- _class: ppt -->", ""]
    if slide.title:
        out.append(f'<div class="page-title">{esc(slide.title)}</div>')
        out.append("")
    out.append('<div class="ppt-canvas">')
    for el in slide.elements:
        if el["kind"] == "text":
            out.append(render_text_box(el))
        elif el["kind"] == "image":
            out.append(render_image(el))
        elif el["kind"] == "table":
            out.append(render_table(el))
        else:
            out.append(render_ph(el))
    out.append("</div>")
    if slide.notes:
        out.append("")
        out.append(f"<!-- {slide.notes.strip().replace('-->', '--&gt;')} -->")
    return "\n".join(out).rstrip() + "\n"


STYLE_BLOCK = """style: |
  /* === pptx2marp: 座標再現レイヤ（色・字体は academic CSS に委ねる） === */
  section.ppt { padding: 0; }
  section.ppt .ppt-canvas { position: absolute; inset: 0; }
  section.ppt .ppt-box {
    position: absolute; box-sizing: border-box; overflow: hidden;
    display: flex; flex-direction: column; justify-content: flex-start;
    padding: 1px 2px;
  }
  section.ppt .ppt-p { margin: 0; line-height: 1.25; }
  section.ppt .ppt-bu { margin-right: .4em; }
  section.ppt .ppt-img { position: absolute; object-fit: contain; }
  section.ppt .ppt-tbl { position: absolute; border-collapse: collapse; font-size: 18px; }
  section.ppt .ppt-tbl th, section.ppt .ppt-tbl td {
    border: 1px solid #bbb; padding: 2px 6px; vertical-align: middle;
  }
  section.ppt .ppt-tbl th { background: var(--accent-soft, #f3e6e8); }
  section.ppt .ppt-ph {
    position: absolute; box-sizing: border-box; display: flex;
    align-items: center; justify-content: center; text-align: center;
    border: 2px dashed var(--accent, #C0182B); color: var(--accent, #C0182B);
    font-size: 14px; padding: 4px;
  }
"""


def build_frontmatter(deck_title: str, ratio: float) -> str:
    size = "16:9" if ratio > 1.5 else "4:3"
    header = (f'<div class="hdr-left">{esc(deck_title)}</div>'
              f'<img class="hdr-logo" src="../assets/chiba-logo.png" alt="CHIBA UNIVERSITY">')
    return (
        "---\n"
        "marp: true\n"
        "theme: academic\n"
        "paginate: true\n"
        f"size: {size}\n"
        f"header: '{header}'\n"
        "footer: ''\n"
        "html: true\n"
        + STYLE_BLOCK +
        "---\n"
    )


# --------------------------------------------------------------------------- #
# main
# --------------------------------------------------------------------------- #
def main():
    ap = argparse.ArgumentParser(description="PowerPoint(.pptx) → 座標再現 + academic テーマの Marp デック")
    ap.add_argument("input", help="入力 .pptx のパス")
    ap.add_argument("--name", help="デック名（YYYYMMDD_<略号>_<内容>）。省略時は自動生成")
    ap.add_argument("--title", help="ヘッダー帯のデックタイトル。省略時はファイル名")
    ap.add_argument("--event", default="Import", help="自動デック名のイベント略号（既定 Import）")
    ap.add_argument("--date", help="YYYYMMDD（省略時は今日）")
    ap.add_argument("--slug", help="内容スラッグ(ascii)。省略時はファイル名から")
    ap.add_argument("--out", default="slides", help="出力ベース（既定 slides）")
    args = ap.parse_args()

    in_path = Path(args.input).expanduser()
    if not in_path.exists():
        sys.exit(f"入力が見つかりません: {in_path}")
    if in_path.suffix.lower() != ".pptx":
        sys.exit("このツールは .pptx 専用です（.pdf は未対応）。")

    date = args.date or _dt.date.today().strftime("%Y%m%d")
    slug = args.slug or slugify(in_path.stem) or "import"
    deck_name = args.name or f"{date}_{args.event}_{slug}"

    base = (REPO / args.out) if not Path(args.out).is_absolute() else Path(args.out)
    deck_dir = base / deck_name
    if deck_dir.exists():
        sys.exit(f"既に存在します: {deck_dir}（--name で別名を）")
    src_dir = deck_dir / "src"
    src_dir.mkdir(parents=True)
    (deck_dir / "out").mkdir()

    prs = Presentation(str(in_path))
    W, H = prs.slide_width, prs.slide_height
    ratio = W / H
    scale = MARP_H_PX * EMU_PER_PT / H  # pt → px（720px キャンバス基準）

    fig = [0]
    identity = (lambda x, y, w, h: (x, y, w, h))
    slides_md, n_unsupported, n_media, n_chart, n_table = [], 0, 0, 0, 0
    for s in prs.slides:
        slide = Slide()
        ctx = {"W": W, "H": H, "src": src_dir, "fig": fig, "scale": scale}
        walk(s.shapes, identity, slide, ctx, s.shapes.title, lift_title=not is_title_slide(s))
        if s.has_notes_slide and s.notes_slide.notes_text_frame:
            slide.notes = s.notes_slide.notes_text_frame.text.strip()
        n_unsupported += sum(1 for e in slide.elements if e["kind"] == "image" and not e["ok"])
        n_media += sum(1 for e in slide.elements if e["kind"] == "ph" and "動画" in e.get("label", ""))
        n_chart += sum(1 for e in slide.elements if e["kind"] == "ph" and "グラフ" in e.get("label", ""))
        n_table += sum(1 for e in slide.elements if e["kind"] == "table")
        slides_md.append(render_slide(slide))

    deck_title = args.title or in_path.stem

    if THEME_CSS.exists():
        shutil.copy2(THEME_CSS, deck_dir / "academic.css")

    md = build_frontmatter(deck_title, ratio) + "\n" + "\n---\n\n".join(slides_md)
    md_path = deck_dir / f"{deck_name}.md"
    md_path.write_text(md, encoding="utf-8")

    rel = md_path.relative_to(REPO)
    print(f"✅ 変換完了: {len(prs.slides)} スライド / 画像 {fig[0]} 枚")
    print(f"   出力: {deck_dir.relative_to(REPO)}/  ({md_path.name} + src/ + out/)")
    print()
    print("プレビュー / PDF（リポジトリルートで）:")
    print(f'  npx @marp-team/marp-cli@latest "{rel}" --theme-set theme/academic.css --html --preview')
    print(f'  tools/marp-pdf/build-pdf.sh "{rel}"')
    if n_table:
        print(f"   表 {n_table} 個を HTML テーブルで再現しました。")
    if n_unsupported:
        print(f"\n⚠ WMF/EMF 画像 {n_unsupported} 枚は描画されません（点線枠で表示）。PowerPoint で PNG 等に再保存を。")
    if n_media:
        print(f"⚠ 動画 {n_media} 件は埋め込めません（位置に注記を表示）。")
    if n_chart:
        print(f"⚠ グラフ {n_chart} 個は再現不可（位置に注記）。marp-echarts で作り直すと綺麗です。")


if __name__ == "__main__":
    main()
